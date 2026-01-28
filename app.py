import streamlit as st
import pandas as pd
import datetime
import hashlib
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches, Pt
import io

# 1. 앱 페이지 설정 (동양의 미와 현대적 세련미의 조화)
st.set_page_config(page_title="황산스님 : 천기비결", page_icon="🏮", layout="wide")

st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Nanum+Myeongjo:wght@400;700&display=swap');
    .main { background-color: #0a0a0c; color: #e0c068; font-family: 'Nanum Myeongjo', serif; }
    .stButton>button { width: 100%; background-color: #b38b2d; color: #fff; font-weight: bold; border-radius: 0px; height: 4em; border: 1px solid #e0c068; font-size: 1.2em; }
    .report-card { background-color: #1a1a1d; padding: 40px; border: 1px solid #3d3d3d; border-radius: 5px; line-height: 2.2; color: #d1d1d1; margin-bottom: 30px; }
    .pillar-box { background-color: #26262b; border: 2px solid #e0c068; padding: 15px; text-align: center; color: #e0c068; font-weight: bold; font-size: 1.5em; }
    h1, h2, h3 { color: #e0c068; text-align: center; }
    .highlight { color: #f1c40f; font-weight: bold; }
    </style>
    """, unsafe_allow_html=True)

# 2. 명리학 분석 엔진 (중국/일본/한국 로직 통합)
class GrandMasterEngine:
    def __init__(self, name, birth, time_str):
        self.name = name
        self.birth = birth
        self.time = time_str
        # 고유 해시 생성으로 수천 가지 조합 구현
        self.seed = int(hashlib.sha256(f"{name}{birth}{time_str}".encode()).hexdigest(), 16)
        
    def get_pillars(self):
        # 사주 팔자(4주 8자) 시뮬레이션
        stems = ["甲", "乙", "丙", "丁", "戊", "己", "庚", "辛", "壬", "癸"]
        branches = ["子", "丑", "寅", "卯", "辰", "巳", "午", "未", "申", "酉", "戌", "亥"]
        return [
            (stems[self.seed % 10], branches[self.seed % 12]), # 시주
            (stems[(self.seed//10) % 10], branches[(self.seed//12) % 12]), # 일주
            (stems[(self.seed//100) % 10], branches[(self.seed//144) % 12]), # 월주
            (stems[(self.seed//1000) % 10], branches[(self.seed//1728) % 12]) # 년주
        ]

    def get_analysis(self):
        # 오행 분석 및 인생 단계별 심층 텍스트
        elements = ["木 (창조, 성취)", "火 (열정, 확산)", "土 (신뢰, 중재)", "金 (결단, 재물)", "水 (지혜, 흐름)"]
        yong_sin = elements[self.seed % 5]
        
        early = [
            "초년운(1~25세): <span class='highlight'>문창귀인(文昌貴人)</span>이 임하여 영특함이 남다릅니다. 봄날의 대지에 내리는 비처럼 부모의 전폭적인 지지 아래 학문적 성취가 높으며, 일찍이 자신의 천직을 발견하여 기초를 탄탄히 다지는 시기입니다. 다만 사주 내 비겁(比劫)이 강해 친구로 인한 손재수가 있으니 대인관계에 신중함이 필요합니다.",
            "초년운(1~25세): <span class='highlight'>식신생재(食神生財)</span>의 격국을 이루어 재능이 곧 재물로 이어지는 운세입니다. 어린 시절부터 예술적 감각이나 기술적 재능이 뛰어나 주변의 찬사를 받습니다. 20대 초반에 강력한 역마운이 들어오니 고향을 떠나 타지에서 공부하거나 활동할 때 운의 크기가 수배로 커지는 흐름을 보입니다."
        ]
        
        mid = [
            "중년운(26~55세): 인생의 황금기인 <span class='highlight'>정관(正官)과 정인(正印)</span>이 상생하는 흐름입니다. 사회적 지위가 급격히 상승하며 자신의 분야에서 일가를 이루게 됩니다. 특히 40대 중반에 천을귀인(天乙貴人)의 조력으로 거대한 문서운이 들어오니 부동산이나 큰 계약을 통해 노후 자금의 기틀을 마련하게 됩니다. 명예와 실리를 동시에 거머쥐는 시기입니다.",
            "중년운(26~55세): <span class='highlight'>편재격(偏財格)</span>이 발동하여 사업적 수완이 극에 달하는 시기입니다. 일반적인 월급 생활보다는 자신의 사업이나 투자를 통해 큰 부를 축적하는 기운이 강합니다. 특히 해외와의 인연이 깊어 구매대행이나 유통업에서 큰 두각을 나타내며, 사람을 다루는 통솔력이 빛을 발해 만인의 우두머리가 되는 형국입니다."
        ]
        
        late = [
            "말년운(56세 이후): <span class='highlight'>식신(食神)</span>이 노년까지 건재하니 자손이 번창하고 신체가 강건합니다. 창고에 곡식이 가득 찬 형국으로 베푸는 삶을 살게 되며, 후학을 양성하거나 사회적 멘토로서 명성을 떨치게 됩니다. 태평성대의 기운이 집안을 감싸니 근심 걱정 없는 안락한 황혼을 보내게 되는 대기만성형의 표본입니다.",
            "말년운(56세 이후): <span class='highlight'>천수성(天壽星)</span>이 길하게 작용하여 무병장수하며, 산 좋고 물 맑은 곳에서 여유를 즐기는 삶이 보입니다. 젊은 시절 쌓아온 인덕이 보답으로 돌아와 귀인들이 끊이지 않으며, 명예로운 직함을 유지하며 품격 있는 노후를 보내게 됩니다. 정신적 지주로서 많은 이들에게 지혜를 전수하는 고귀한 삶입니다."
        ]

        return {
            "yong_sin": yong_sin,
            "early": early[self.seed % 2],
            "mid": mid[(self.seed // 2) % 2],
            "late": late[(self.seed // 4) % 2],
            "balance": [self.seed%30, (self.seed//7)%25, (self.seed//13)%20, (self.seed//19)%35, (self.seed//3)%15]
        }

# 3. 메인 화면
st.markdown("<h1 style='font-size: 3.5em;'>🏮 황산스님 천기비결(天機秘訣)</h1>", unsafe_allow_html=True)
st.markdown("<p style='text-align:center; font-size:1.2em;'>중국·일본·한국 3국 명리학의 정수를 집약한 글로벌 하이엔드 운명 분석</p>", unsafe_allow_html=True)

with st.container():
    st.markdown("<div class='report-card'>", unsafe_allow_html=True)
    c1, c2, c3, c4 = st.columns(4)
    with c1: name = st.text_input("고객 성함", "방문객")
    with c2: birth = st.date_input("생년월일", datetime.date(1980, 1, 1))
    with c3: time_str = st.selectbox("출생 시간", [f"{i:02d}시" for i in range(24)])
    with c4: gender = st.radio("성별", ["남성", "여성"], horizontal=True)
    
    if st.button("🔮 황산스님의 천기(天機) 분석 시작"):
        engine = GrandMasterEngine(name, birth, time_str)
        pillars = engine.get_pillars()
        res = engine.get_analysis()
        
        st.divider()
        st.markdown(f"### ✨ {name}님의 사주팔자(四柱八字) 원국")
        
        p1, p2, p3, p4 = st.columns(4)
        labels = ["시주(時柱)", "일주(日柱)", "월주(月柱)", "년주(年柱)"]
        for i, col in enumerate([p1, p2, p3, p4]):
            with col:
                st.markdown(f"<div style='text-align:center; color:#888;'>{labels[i]}</div>", unsafe_allow_html=True)
                st.markdown(f"<div class='pillar-box'>{pillars[i][0]}<br>{pillars[i][1]}</div>", unsafe_allow_html=True)

        st.divider()
        
        # 오행 분포 그래프
        st.markdown("### 📊 오행(五行) 에너지 균형도")
        fig, ax = plt.subplots(figsize=(10, 2), facecolor='#1a1a1d')
        ax.set_facecolor('#1a1a1d')
        bars = ax.barh(['木', '火', '土', '金', '水'], res['balance'], color=['#27ae60', '#e74c3c', '#f1c40f', '#ecf0f1', '#2980b9'])
        ax.tick_params(colors='#e0c068')
        st.pyplot(fig)
        
        st.markdown(f"<p style='text-align:center;'>당신의 삶을 이끄는 핵심 기운(용신)은 <span class='highlight'>{res['yong_sin']}</span>입니다.</p>", unsafe_allow_html=True)

        # 심층 분석 탭
        tabs = st.tabs(["🌱 초년/성장", "☀️ 중년/성취", "🌕 말년/안식", "💰 재물/직업", "🧘 황산스님의 비책"])
        
        with tabs[0]:
            st.markdown(f"<div class='report-card'>{res['early']}</div>", unsafe_allow_html=True)
            
        with tabs[1]:
            st.markdown(f"<div class='report-card'>{res['mid']}</div>", unsafe_allow_html=True)
            
        with tabs[2]:
            st.markdown(f"<div class='report-card'>{res['late']}</div>", unsafe_allow_html=True)
            
        with tabs[3]:
            st.markdown("<div class='report-card'><b>[재물운]</b> 정재와 편재가 조화로우니 평생 돈이 마르지 않는 명입니다. 특히 4, 9가 들어간 날짜에 큰 계약이 성사될 확률이 높습니다.<br><b>[직업운]</b> 전문 기술이나 라이선스를 바탕으로 한 고부가가치 산업 혹은 사람을 가르치고 살리는 교육/의료/상담업이 천직입니다.</div>", unsafe_allow_html=True)
        with tabs[4]:
            st.markdown(f"<div class='report-card'><b>'운명은 30%가 정해진 길이고, 70%는 내가 닦아가는 길이다.'</b><br>{name}님, 당신의 사주는 금(金)의 기운이 강하니 서쪽을 가까이하고, 흰색 계열의 옷을 입어 기운을 보강하십시오. 매일 아침 432Hz의 주파수를 들으며 마음을 정화하면 막혔던 재물운이 폭포수처럼 쏟아질 것입니다.</div>", unsafe_allow_html=True)

        # 리포트 다운로드
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"{name}님의 인생 천기비결"
        tf = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5)).text_frame
        tf.word_wrap = True
        tf.text = f"[초년] {res['early'][:100]}...\n\n[중년] {res['mid'][:100]}...\n\n[말년] {res['late'][:100]}..."
        buf = io.BytesIO()
        prs.save(buf)
        st.download_button("📥 5만원 상당의 프리미엄 평생 운세장 소장하기", buf.getvalue(), file_name=f"{name}_황산스님_리포트.pptx")
